#!/usr/bin/env python3
"""
future.nostalgia venue deck → fully-editable PowerPoint.
Ports the verified HTML deck: 1280x720 stage → 16:9 slide (inches=px/96, pt=px*0.75).
Every element is a native textbox/shape. Dieter Rams: flat (no shadows), 0 radius,
1px hairlines, one teal accent + one coral ember per slide; Jost (display) + Inter (text).
CURDY vertically centers each slide's body; the footer stays pinned.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

CANVAS="F4F1EA"; PAPER="FBF8F1"; INK="2D2D2A"; INKM="5A5A55"; INKF="9A988F"
HAIR="D1D1CE"; HAIRS="C5C2B6"; ACCENT="2B5B5D"; EMBER="FF5A36"; PHOTO="CDD4CE"
JOST="Jost"; INTER="Inter"
ASSET="/Users/josh/developer/futurenostalgia.ca/assets/brand"
CENTER=PP_ALIGN.CENTER; RIGHTA=PP_ALIGN.RIGHT; LEFTA=PP_ALIGN.LEFT
CURDY=0  # per-slide vertical shift for body content (footer excluded)

def IN(px): return Inches(px/96.0)
def PTpx(px): return Pt(px*0.75)

prs = Presentation()
prs.slide_width = Emu(int(13.3333*914400)); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor.from_string(CANVAS)
    return s

def noshadow(sp):
    """Flatten: strip the theme effectRef (a:style) that renders as a drop shadow."""
    el = sp._element
    for st in el.findall(qn('p:style')):
        el.remove(st)
    try: sp.shadow.inherit = False
    except Exception: pass

def rect(s, x, y, w, h, fill=None, line=None, line_px=1, rot=0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, IN(x), IN(y+CURDY), IN(w), IN(h))
    noshadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = PTpx(line_px)
    if rot: sp.rotation = rot
    return sp

def hline(s, x, y, w, color=HAIR, px=1): return rect(s, x, y, w, px, fill=color)
def vline(s, x, y, h, color=HAIR, px=1): return rect(s, x, y, px, h, fill=color)

def set_track(run, px): run._r.get_or_add_rPr().set('spc', str(int(round(px*0.75*100))))

def txt(s, x, y, w, h, segs, size, color=INK, font=INTER, bold=False,
        align=LEFTA, line=1.0, track=0, upper=False, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(IN(x), IN(y+CURDY), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    if line: p.line_spacing = line
    if isinstance(segs, str): segs = [(segs, {})]
    for t, o in segs:
        r = p.add_run(); r.text = t.upper() if upper else t
        r.font.size = PTpx(o.get("size", size)); r.font.name = o.get("font", font)
        r.font.bold = o.get("bold", bold)
        r.font.color.rgb = RGBColor.from_string(o.get("color", color))
        tr = o.get("track", track)
        if tr: set_track(r, tr)
    return tb

def arrow(s, x1, y1, x2, y2, color=INKF, w=1.5):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1+CURDY), IN(x2), IN(y2+CURDY))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = PTpx(w)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type':'triangle','w':'med','len':'med'}))
    noshadow(c); return c

def pic(s, path, x, y, h=None, w=None, rot=0):
    kw = {}
    if h: kw["height"] = IN(h)
    if w: kw["width"] = IN(w)
    p = s.shapes.add_picture(path, IN(x), IN(y+CURDY), **kw)
    if rot: p.rotation = rot
    return p

LEFT=84; RIGHT=1196; CW=1112; FOOT_HAIR=648; FOOT_Y=657

def kicker(s, x, y, label, color=ACCENT):
    rect(s, x, y, 44, 2, fill=color)
    txt(s, x, y+14, 600, 20, label, 13, color=color, bold=True, track=1.6, upper=True)

def foot(s, num):
    global CURDY
    save=CURDY; CURDY=0
    hline(s, LEFT, FOOT_HAIR, CW, HAIR)
    pic(s, f"{ASSET}/wordmark-black.png", LEFT, FOOT_Y, h=16)
    txt(s, RIGHT-160, FOOT_Y-1, 160, 20,
        [(num, {"color":ACCENT}), (" / 06", {"color":INKM})], 14, font=JOST, bold=True, align=RIGHTA)
    CURDY=save

def print_frame(s, x, y, w, photo_h, cap, tag=None, reg=True, mat_bottom=40, pad=13):
    ph_w = w - 2*pad
    rect(s, x, y, w, pad+photo_h+mat_bottom, fill=PAPER, line=HAIRS, line_px=1)
    rect(s, x+pad, y+pad, ph_w, photo_h, fill=PHOTO)
    if tag:
        tw = len(tag)*6.8 + 18
        rect(s, x+pad+8, y+pad+photo_h-26, tw, 18, fill=CANVAS, line=ACCENT, line_px=1)
        txt(s, x+pad+8, y+pad+photo_h-23, tw, 14, tag, 10, color=ACCENT, bold=True, track=0.8,
            upper=True, align=CENTER)
    cy = y+pad+photo_h+(mat_bottom-18)/2
    txt(s, x+pad, cy, ph_w-22, 16, cap, 12, color=INKM)
    if reg:
        rect(s, x+w-pad-10, cy+1, 10, 10, line=ACCENT, line_px=1.5, shape=MSO_SHAPE.OVAL)

# ════════════ SLIDE 1 — COVER ════════════
CURDY=0; s = slide()
pic(s, f"{ASSET}/wordmark-black.png", LEFT, 138, h=46)
txt(s, LEFT, 208, 600, 20, "Analog photo booth for events", 13, color=ACCENT, bold=True, track=1.6, upper=True)
txt(s, LEFT, 244, 560, 240, "Photos your guests actually keep.", 74, font=JOST, bold=True, line=0.98, track=-1.4)
txt(s, LEFT, 500, 575, 120,
    "I bring an analog print booth to your event and run the whole night. Every guest walks out with a real photograph in their hand.",
    20, color=INKM, line=1.5)
txt(s, LEFT, 614, 150, 20, "futurenostalgia.ca", 13, color=INKM)
vline(s, LEFT+158, 616, 16, HAIR)
txt(s, LEFT+174, 614, 300, 20, "Festivals, hotels & events", 13, color=INKM)
print_frame(s, 856, 300, 340, 209, "Harbourfront · 11:48pm", tag="Your event, here")
foot(s, "01")

# ════════════ SLIDE 2 — WHY PRINTS ════════════
CURDY=-85; s = slide()
PHX, PHY, PHW, PHH = LEFT, 252, 470, 352
rect(s, PHX, PHY, PHW, PHH, fill=PHOTO, line=HAIRS, line_px=1)
rect(s, PHX+12, PHY+12, 14, 2, fill=ACCENT); rect(s, PHX+12, PHY+12, 2, 14, fill=ACCENT)
cyc = PHY+PHH/2
txt(s, PHX, cyc-20, PHW, 18, "ADD A PHOTO", 11, color=INKF, bold=True, track=1.4, upper=True, align=CENTER)
txt(s, PHX, cyc+2, PHW, 18, "a guest holding their print", 11, color=INKF, align=CENTER)
RX = 600
kicker(s, RX, 252, "Why prints")
txt(s, RX, 300, 580, 200,
    [("We don't remember files. We remember ", {}), ("objects.", {"color":EMBER})],
    58, font=JOST, bold=True, line=1.06, track=-0.8)
txt(s, RX, 532, 470, 120,
    "Your guests will shoot a hundred photos tonight and never open them again. The one they can hold, they keep for years. A print is a memory with a body.",
    18, color=INKM, line=1.5)
foot(s, "02")

# ════════════ SLIDE 3 — HOW IT WORKS ════════════
CURDY=-70; s = slide()
kicker(s, LEFT, 210, "How it works")
txt(s, LEFT, 250, 700, 70, "I run the whole booth.", 46, font=JOST, bold=True, track=-0.7)
DT = 350
def node(x): return LEFT + x
ny = DT + 20
for bx in (8, 268, 528, 788):
    rect(s, node(bx), ny, 104, 104, line=HAIRS, line_px=1.5)
cyc = ny + 52
for x1, x2 in ((120, 260), (380, 520), (640, 780)):
    arrow(s, node(x1), cyc, node(x2), cyc, color=INKF, w=1.5)
rect(s, node(40), ny+32, 2.5, 40, fill=ACCENT); rect(s, node(40), ny+70, 40, 2.5, fill=ACCENT)
rect(s, node(296), ny+38, 48, 34, line=ACCENT, line_px=2.5)
rect(s, node(311), ny+47, 18, 18, line=ACCENT, line_px=2.5, shape=MSO_SHAPE.OVAL)
rect(s, node(305), ny+30, 22, 8, line=ACCENT, line_px=2.5)
rect(s, node(566), ny+20, 28, 22, line=ACCENT, line_px=2.5)
rect(s, node(556), ny+54, 48, 30, line=ACCENT, line_px=2.5)
rect(s, node(816), ny+24, 48, 56, line=EMBER, line_px=2.5, rot=-8)
labs = [("Your corner","a few sq. ft. + power",8),("I shoot it","camera + lighting, all mine",268),
        ("Prints on the spot","real film, about a minute",528),("They keep it","your event on the border",788)]
for name, sub, bx in labs:
    txt(s, node(bx), DT+150, 200, 18, name, 13, bold=True, track=1.3, upper=True)
    txt(s, node(bx), DT+170, 220, 18, sub, 12.5, color=INKM)
txt(s, node(940), DT+50, 170, 80, "Your staff never touches any of it.", 13, color=INKM, line=1.35)
txt(s, LEFT, 578, 900, 56,
    "You give me a corner with power. I bring everything else and pack it all down at the end.",
    18, color=INKM, line=1.5)
foot(s, "03")

# ════════════ SLIDE 4 — WHAT IT COSTS YOU ════════════
CURDY=-65; s = slide()
kicker(s, LEFT, 198, "What it costs you")
txt(s, LEFT, 232, 700, 120, "Nothing.", 92, font=JOST, bold=True, track=-1.4)
txt(s, LEFT, 360, 520, 60, "Your guests or a sponsor cover the prints. There are two ways to run it.",
    19, color=INKM, line=1.5)
MY = 458
vline(s, 648, MY, 150, HAIR)
rect(s, LEFT, MY, 36, 2, fill=ACCENT)
txt(s, LEFT, MY+16, 460, 36, "Per print", 25, bold=True, track=-0.3)
txt(s, LEFT, MY+56, 470, 90, "Guests pay for their own prints on the spot, the way they'd buy anything else at your event.",
    16, color=INKM, line=1.55)
RX = 708
rect(s, RX, MY, 36, 2, fill=EMBER)
txt(s, RX, MY+16, 460, 36, "Sponsored", 25, bold=True, track=-0.3)
txt(s, RX, MY+56, 470, 90, "A sponsor covers a set number, so guests take them home free, with the sponsor's name on every print.",
    16, color=INKM, line=1.55)
txt(s, LEFT, 618, 800, 20, "I'll help you pick the one that fits your event.", 14, color=INKM)
foot(s, "04")

# ════════════ SLIDE 5 — WHAT THEY TAKE HOME ════════════
CURDY=-60; s = slide()
kicker(s, LEFT, 210, "What they take home")
txt(s, LEFT, 250, 800, 110, "A real print. Your event on the border.", 46, font=JOST, bold=True, line=1.04, track=-0.7)
ROWY = 372; pw = (CW - 3*16)/4
cards = [("Front row · 9:42pm","Your event",True),("The line, all night",None,False),
         ("Two strangers, now friends",None,False),("Last dance · 1:10am","Your event",True)]
for i,(cap,tag,reg) in enumerate(cards):
    print_frame(s, LEFT + i*(pw+16), ROWY, pw, (pw-18)*2/3, cap, tag=tag, reg=reg, mat_bottom=30, pad=9)
txt(s, LEFT, 596, 800, 20, "Swap these for shots from your room, your crowd, your night.", 14, color=INKM)
foot(s, "05")

# ════════════ SLIDE 6 — LET'S TALK + context ════════════
CURDY=0; s = slide()
kicker(s, LEFT, 120, "Let's talk")
txt(s, LEFT, 166, 840, 170,
    [("Give me a corner. I'll handle the rest. ", {}), ("→", {"color":EMBER})],
    64, font=JOST, bold=True, line=1.04, track=-1.2)
txt(s, LEFT, 320, 820, 60,
    "Tell me about your event and I'll come back within a day with dates and a price. No deposit to ask.",
    20, color=INKM, line=1.5)
cy = 402
lx = LEFT
links = ["hello@futurenostalgia.ca", "futurenostalgia.ca", "@future.nostalgia"]
for i, lk in enumerate(links):
    lw = 24 + len(lk)*8.4
    txt(s, lx, cy, lw, 24, lk, 16, color=INK); lx += lw
    if i < len(links)-1:
        vline(s, lx+8, cy+2, 18, HAIR); lx += 40
hline(s, LEFT, 452, CW, HAIR)
pic(s, f"{ASSET}/about-family.png", LEFT, 474, w=104, rot=-3)
txt(s, 214, 480, 720, 90,
    [("I'm Josh Gonzales.", {"bold":True}),
     (" I started future.nostalgia for my own family, because the photos we keep are the ones we can hold. A share of every event goes to brain research — the science of why a photograph is worth keeping at all.", {})],
    14, color=INKM, line=1.55)
foot(s, "06")

OUT = "/Users/josh/developer/futurenostalgia.ca/deck/future-nostalgia-deck.pptx"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
